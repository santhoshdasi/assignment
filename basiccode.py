from pptx import Presentation
from pptx.util import Inches

# Create a new PowerPoint presentation
presentation = Presentation()

# Load the font file
font_file = 'sample_font_file.ttf'

# Read content from sample_slide1_input.txt
with open('sample_slide1_input.txt', 'r') as slide1_file:
    slide1_content = slide1_file.read()

# Create Slide 1
slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
title1 = slide1.shapes.title
title1.text = "Slide 1"

# Add the content to Slide 1
content1 = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
text_frame1 = content1.text_frame
p1 = text_frame1.add_paragraph()
p1.text = slide1_content

# Apply the font style to Slide 1
for paragraph in text_frame1.paragraphs:
    for run in paragraph.runs:
        run.font.file = font_file
        run.font.name = run.font.file.name

# Read content from sample_slide2_input.txt
with open('sample_slide2_input.txt', 'r') as slide2_file:
    slide2_content = slide2_file.read()

# Create Slide 2
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title2 = slide2.shapes.title
title2.text = "Slide 2"

# Add the content to Slide 2
content2 = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
text_frame2 = content2.text_frame
p2 = text_frame2.add_paragraph()
p2.text = slide2_content

# Apply the font style to Slide 2
for paragraph in text_frame2.paragraphs:
    for run in paragraph.runs:
        run.font.file = font_file
        run.font.name = run.font.file.name

# Save the presentation
presentation.save("output.pptx")
