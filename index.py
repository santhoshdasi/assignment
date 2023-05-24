from pptx import Presentation
from pptx.util import Inches
presentation = Presentation()

slide_layout = presentation.slide_layouts[1]  

slide = presentation.slides.add_slide(slide_layout)  
title = slide.shapes.title  
title.text = "My First Slide"  
content = slide.placeholders[1] 
content.text = "Lorem ipsum dolor sit amet, consectetur adipiscing elit."   
slide_layout = presentation.slide_layouts[1]  

slide = presentation.slides.add_slide(slide_layout)  
title = slide.shapes.title  
title.text = "My second Slide"  
content = slide.placeholders[1] 
content.text = "Lorem ipsum sit afoyt met, hoihebfcfwef adipiscing elit."   



presentation.save("my_presentation.pptx")
